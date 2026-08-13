import io
import os
from typing import List, Dict, Any

class DocumentLoader:
    """Multi-format document loader supporting PDF, DOCX, PPTX, and TXT files."""
    
    @staticmethod
    def load_txt(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        text = file_content.decode("utf-8", errors="ignore")
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
        
        pages = []
        # Group into pseudo-pages or sections
        chunk_size = 5
        for i in range(0, len(paragraphs), chunk_size):
            page_text = "\n\n".join(paragraphs[i:i+chunk_size])
            pages.append({
                "text": page_text,
                "metadata": {
                    "source": filename,
                    "page": (i // chunk_size) + 1,
                    "file_type": "TXT"
                }
            })
        if not pages and text.strip():
            pages.append({
                "text": text.strip(),
                "metadata": {"source": filename, "page": 1, "file_type": "TXT"}
            })
        return pages

    @staticmethod
    def load_pdf(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(file_content))
            pages = []
            for idx, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    pages.append({
                        "text": page_text.strip(),
                        "metadata": {
                            "source": filename,
                            "page": idx + 1,
                            "file_type": "PDF"
                        }
                    })
            return pages
        except Exception as e:
            print(f"Error reading PDF {filename}: {e}")
            return []

    @staticmethod
    def load_docx(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_content))
            full_text = []
            for p in doc.paragraphs:
                if p.text.strip():
                    full_text.append(p.text.strip())
            
            combined_text = "\n\n".join(full_text)
            # Group into sections as pseudo-pages
            sections = []
            paragraphs = [p for p in full_text if p]
            batch_size = 6
            for i in range(0, len(paragraphs), batch_size):
                section_text = "\n".join(paragraphs[i:i+batch_size])
                sections.append({
                    "text": section_text,
                    "metadata": {
                        "source": filename,
                        "page": (i // batch_size) + 1,
                        "file_type": "DOCX"
                    }
                })
            return sections if sections else [{"text": combined_text, "metadata": {"source": filename, "page": 1, "file_type": "DOCX"}}]
        except Exception as e:
            print(f"Error reading DOCX {filename}: {e}")
            return []

    @staticmethod
    def load_pptx(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_content))
            slides = []
            for idx, slide in enumerate(prs.slides):
                slide_text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text_parts.append(shape.text.strip())
                slide_text = "\n".join(slide_text_parts)
                if slide_text.strip():
                    slides.append({
                        "text": slide_text.strip(),
                        "metadata": {
                            "source": filename,
                            "page": idx + 1,
                            "file_type": "PPTX"
                        }
                    })
            return slides
        except Exception as e:
            print(f"Error reading PPTX {filename}: {e}")
            return []

    @classmethod
    def load_file(cls, file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
        ext = os.path.splitext(filename)[1].lower()
        if ext == ".pdf":
            return cls.load_pdf(file_bytes, filename)
        elif ext in [".docx", ".doc"]:
            return cls.load_docx(file_bytes, filename)
        elif ext in [".pptx", ".ppt"]:
            return cls.load_pptx(file_bytes, filename)
        else:
            return cls.load_txt(file_bytes, filename)
