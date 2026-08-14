import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_path="AI_Tutor_Capstone_Presentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]

    # Non-Blue Color Palette (Emerald Teal & Warm Amber Theme)
    COLOR_BG = RGBColor(13, 17, 23)          # #0D1117 Dark Charcoal
    COLOR_CARD = RGBColor(22, 27, 34)        # #161B22 Deep Slate Card
    COLOR_BORDER = RGBColor(33, 38, 45)      # #21262D Border
    COLOR_ACCENT = RGBColor(16, 185, 129)    # #10B981 Emerald Teal Primary Accent
    COLOR_ACCENT_ALT = RGBColor(245, 158, 11) # #F59E0B Warm Amber Secondary Accent
    COLOR_TEXT_MAIN = RGBColor(240, 246, 252) # #F0F6FC Crisp Text
    COLOR_TEXT_MUTED = RGBColor(139, 148, 158) # #8B949E Muted Text
    COLOR_EMERALD_LIGHT = RGBColor(52, 211, 153) # #34D399 Light Mint Emerald

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, title_text, category_text="CAPSTONE PROJECT PRESENTATION"):
        # Header Badge / Category
        badge_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.4))
        tf_badge = badge_box.text_frame
        tf_badge.word_wrap = True
        p_badge = tf_badge.paragraphs[0]
        p_badge.text = category_text.upper()
        p_badge.font.size = Pt(10)
        p_badge.font.bold = True
        p_badge.font.color.rgb = COLOR_EMERALD_LIGHT
        p_badge.font.name = "Arial"

        # Slide Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_MAIN
        p_title.font.name = "Arial"

    def add_footer(slide, current_slide, total_slides=10):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.733), Inches(0.4))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"AI-Based Personalized Tutor Using RAG for Syllabus-Aligned Learning  |  Slide {current_slide} of {total_slides}"
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_TEXT_MUTED
        p.font.name = "Arial"

    def create_card(slide, left, top, width, height, title="", border_accent=True, accent_color=COLOR_ACCENT):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = COLOR_BORDER
        shape.line.width = Pt(1)

        if border_accent:
            accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06))
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = accent_color
            accent_bar.line.fill.background()

        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.5))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = COLOR_EMERALD_LIGHT
            p.font.name = "Arial"

        return shape

    # ==========================================
    # SLIDE 1: Title Slide (Emerald & Gold Theme)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)

    create_card(slide1, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1), border_accent=True, accent_color=COLOR_ACCENT)

    t_box = slide1.shapes.add_textbox(Inches(1.8), Inches(1.6), Inches(9.733), Inches(4.3))
    tf1 = t_box.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "CAPSTONE PROJECT PRESENTATION"
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_EMERALD_LIGHT
    p1.space_after = Pt(14)

    p2 = tf1.add_paragraph()
    p2.text = "AI-Based Personalized Tutor Using RAG\nfor Syllabus-Aligned Learning"
    p2.font.size = Pt(30)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_MAIN
    p2.space_after = Pt(20)

    p3 = tf1.add_paragraph()
    p3.text = "An Intelligent Educational AI System Grounded Strictly in Course Materials to Eliminate LLM Hallucination & Optimize Exam Preparation"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.space_after = Pt(30)

    p4 = tf1.add_paragraph()
    p4.text = "Technology Stack: Python 3.14 | Streamlit | TF-IDF Vector Database | Google Gemini API"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_ACCENT_ALT

    add_footer(slide1, 1)

    # ==========================================
    # SLIDE 2: Introduction & Problem Statement
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "Introduction & Problem Statement")

    create_card(slide2, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.9), "Problem Statement", accent_color=COLOR_ACCENT)
    tb = slide2.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    bullets1 = [
        ("Exam Preparation Struggles", "Students face immense difficulties preparing for exams due to fragmented resources and lack of targeted study material."),
        ("Generic LLM Hallucinations", "Standard LLMs (ChatGPT, Gemini, Claude) generate generic, out-of-syllabus responses that introduce irrelevant topics."),
        ("Inconsistency & Time Waste", "Reading entire textbooks before exams is extremely time-consuming, while un-grounded AI chatbots waste time with false info."),
        ("Lack of Structured Planning", "Without a personalized prep roadmap aligned with exam timelines, students fail to cover crucial syllabus topics effectively.")
    ]
    for b_title, b_desc in bullets1:
        p = tf.add_paragraph(); p.text = f"•  {b_title}: {b_desc}"; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN; p.space_after = Pt(10)

    create_card(slide2, Inches(6.933), Inches(1.7), Inches(5.6), Inches(4.9), "Problem Impact & Industry Need", accent_color=COLOR_ACCENT_ALT)
    tb2 = slide2.shapes.add_textbox(Inches(7.133), Inches(2.4), Inches(5.2), Inches(4.0))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    bullets2 = [
        ("High Failure & Anxiety Rates", "Students spend hours sifting through non-syllabus content right before exams, leading to poor prep quality."),
        ("No Trust in AI Answers", "Generic AI responses lack source verification, preventing students from cross-checking textbook references."),
        ("Core Requirement", "There is a vital need for an intelligent AI tutor that generates responses strictly from syllabus documents with verifiable page citations."),
        ("Multi-Format Ingestion", "Students need a single platform capable of ingesting PDF textbooks, DOCX notes, PPTX slides, and TXT files seamlessly.")
    ]
    for b_title, b_desc in bullets2:
        p = tf2.add_paragraph(); p.text = f"•  {b_title}: {b_desc}"; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN; p.space_after = Pt(10)

    add_footer(slide2, 2)

    # ==========================================
    # SLIDE 3: Proposed Solution
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "Proposed Solution: RAG-Grounded AI Tutor")

    col_w = Inches(3.644)
    gap = Inches(0.4)

    create_card(slide3, Inches(0.8), Inches(1.7), col_w, Inches(4.9), "Retrieval Augmented Generation", accent_color=COLOR_ACCENT)
    tb = slide3.shapes.add_textbox(Inches(1.0), Inches(2.4), col_w - Inches(0.4), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    items1 = [
        "Converts syllabus documents, lecture notes, and textbooks into vector embeddings stored in a vector database.",
        "Retrieves top-k most relevant syllabus context blocks dynamically for every query.",
        "Prevents hallucinations by restricting LLM generation strictly to retrieved context."
    ]
    for item in items1:
        p = tf.add_paragraph(); p.text = f"• {item}"; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN; p.space_after = Pt(12)

    create_card(slide3, Inches(0.8) + col_w + gap, Inches(1.7), col_w, Inches(4.9), "Verifiable Source Citations", accent_color=COLOR_ACCENT_ALT)
    tb = slide3.shapes.add_textbox(Inches(1.0) + col_w + gap, Inches(2.4), col_w - Inches(0.4), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    items2 = [
        "Every AI response includes exact citations referencing document name, page number, and similarity score.",
        "Students can cross-check AI outputs directly against source textbooks for self-study.",
        "Transparent vector inspection allows testing vector similarity scores in real-time."
    ]
    for item in items2:
        p = tf.add_paragraph(); p.text = f"• {item}"; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN; p.space_after = Pt(12)

    create_card(slide3, Inches(0.8) + (col_w + gap) * 2, Inches(1.7), col_w, Inches(4.9), "Comprehensive Prep Modules", accent_color=COLOR_ACCENT)
    tb = slide3.shapes.add_textbox(Inches(1.0) + (col_w + gap) * 2, Inches(2.4), col_w - Inches(0.4), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    items3 = [
        "Day-by-Day Study Roadmap Generator covering full preparation timeline up to Day N.",
        "Topic-Wise Practice Quiz Builder generating MCQs & short answers with detailed step-by-step solution keys.",
        "Concept Clarifier providing intuitive real-world analogies for difficult topics."
    ]
    for item in items3:
        p = tf.add_paragraph(); p.text = f"• {item}"; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN; p.space_after = Pt(12)

    add_footer(slide3, 3)

    # ==========================================
    # SLIDE 4: System Design & Architecture
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "System Design & Pipeline Architecture")

    create_card(slide4, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.9), "Offline Ingestion Pipeline", accent_color=COLOR_ACCENT)
    tb = slide4.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    pipe1 = [
        ("Multi-Format Document Loader", "Extracts clean raw text from PDF (pypdf), DOCX (python-docx), PPTX (python-pptx), and TXT files."),
        ("Recursive Text Chunker", "Splits raw text into ~450-character overlapping segments with 90-character overlap, preserving metadata tags (source, page, file_type)."),
        ("TF-IDF Vector Space Model", "Transforms text chunks into TF-IDF vector space embeddings with sublinear term frequency scaling."),
        ("In-Memory Vector Database", "Indexes vector embeddings for high-speed local cosine similarity matching and source-level filtering.")
    ]
    for p_title, p_desc in pipe1:
        p = tf.add_paragraph(); p.text = f"1. {p_title}: {p_desc}"; p.font.size = Pt(10); p.font.color.rgb = COLOR_TEXT_MAIN; p.space_after = Pt(8)

    create_card(slide4, Inches(6.933), Inches(1.7), Inches(5.6), Inches(4.9), "Real-Time Query Pipeline", accent_color=COLOR_ACCENT_ALT)
    tb = slide4.shapes.add_textbox(Inches(7.133), Inches(2.4), Inches(5.2), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    pipe2 = [
        ("Student Query Input", "User submits a question or prompt via the interactive Streamlit chat interface."),
        ("Vector Similarity Retriever", "Computes cosine similarity between query vector and database chunks, selecting top-k relevant segments."),
        ("Grounded LLM Generator", "Passes top-k syllabus context blocks + strict system prompt to Google Gemini API (gemini-flash-latest)."),
        ("Citations & Output Rendering", "Renders structured markdown answer with expandable citation cards displaying document name, page, and snippet.")
    ]
    for p_title, p_desc in pipe2:
        p = tf.add_paragraph(); p.text = f"2. {p_title}: {p_desc}"; p.font.size = Pt(10); p.font.color.rgb = COLOR_TEXT_MAIN; p.space_after = Pt(8)

    add_footer(slide4, 4)

    # ==========================================
    # SLIDE 5: Key Features
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "Key System Features")

    card_w2 = Inches(5.6)
    card_h2 = Inches(2.3)

    create_card(slide5, Inches(0.8), Inches(1.7), card_w2, card_h2, "Multi-Syllabus Scope Selector", accent_color=COLOR_ACCENT)
    tb = slide5.shapes.add_textbox(Inches(1.0), Inches(2.3), card_w2 - Inches(0.4), card_h2 - Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    p.text = "Allows students to query across 'All Documents' simultaneously or narrow scope to a specific subject file. Supports duplicate file updates without indexing conflicts."
    p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN

    create_card(slide5, Inches(6.933), Inches(1.7), card_w2, card_h2, "Instant Doubt Resolution with Citations", accent_color=COLOR_ACCENT_ALT)
    tb = slide5.shapes.add_textbox(Inches(7.133), Inches(2.3), card_w2 - Inches(0.4), card_h2 - Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    p.text = "Interactive chat workspace providing instant syllabus-grounded answers. Each response features expandable citation cards with exact document source, page number, and similarity score."
    p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN

    create_card(slide5, Inches(0.8), Inches(4.3), card_w2, card_h2, "Day-by-Day Study Roadmap Generator", accent_color=COLOR_ACCENT_ALT)
    tb = slide5.shapes.add_textbox(Inches(1.0), Inches(4.9), card_w2 - Inches(0.4), card_h2 - Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    p.text = "Builds a complete, un-truncated study schedule spanning from Day 1 to Day N (up to 60 days). Samples topics proportionally across all uploaded syllabus documents."
    p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN

    create_card(slide5, Inches(6.933), Inches(4.3), card_w2, card_h2, "Practice Quiz & Concept Clarifier", accent_color=COLOR_ACCENT)
    tb = slide5.shapes.add_textbox(Inches(7.133), Inches(4.9), card_w2 - Inches(0.4), card_h2 - Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    p.text = "Generates exam-style practice question papers (MCQs + short answer) with detailed solution keys, and clarifies complex topics using intuitive real-world analogies."
    p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN

    add_footer(slide5, 5)

    # ==========================================
    # SLIDE 6: Output / Results (With Screenshot)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6)
    add_header(slide6, "Output & Results — System Interface Dashboard")

    screenshot_path = "dashboard_mockup.png"
    if os.path.exists(screenshot_path):
        slide6.shapes.add_picture(screenshot_path, Inches(0.8), Inches(1.6), Inches(7.8), Inches(4.9))

    create_card(slide6, Inches(8.8), Inches(1.6), Inches(3.733), Inches(4.9), "Results Highlights", accent_color=COLOR_ACCENT)
    tb = slide6.shapes.add_textbox(Inches(9.0), Inches(2.3), Inches(3.333), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    res_items = [
        ("Visual Excellence", "Executive dark charcoal design system with Emerald Teal typography and responsive metric cards."),
        ("Instant Grounded Q&A", "Verified 100% syllabus-aligned responses with collapsible source citation cards."),
        ("Multi-Format Ingestion", "Successfully processes PDF, DOCX, PPTX, and TXT documents."),
        ("Full Day-by-Day Roadmap", "Generates complete, un-truncated prep schedules covering all requested days."),
        ("Zero Database Leakage", "Clean database clearing wipes vector store completely with 0 lingering chunks.")
    ]
    for r_title, r_desc in res_items:
        p = tf.add_paragraph(); p.text = f"• {r_title}: {r_desc}"; p.font.size = Pt(10); p.font.color.rgb = COLOR_TEXT_MAIN; p.space_after = Pt(8)

    add_footer(slide6, 6)

    # ==========================================
    # SLIDE 7: Limitations
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7)
    add_header(slide7, "System Limitations")

    create_card(slide7, Inches(0.8), Inches(1.7), col_w, Inches(4.9), "Input Data Quality Dependency", accent_color=COLOR_ACCENT)
    tb = slide7.shapes.add_textbox(Inches(1.0), Inches(2.4), col_w - Inches(0.4), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    l1 = [
        "RAG output accuracy directly depends on the quality and completeness of uploaded syllabus documents.",
        "Scanned PDFs without OCR or corrupted text files may yield low text extraction quality.",
        "Incomplete syllabus uploads result in partial context retrieval for complex topics."
    ]
    for item in l1:
        p = tf.add_paragraph(); p.text = f"• {item}"; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN; p.space_after = Pt(10)

    create_card(slide7, Inches(0.8) + col_w + gap, Inches(1.7), col_w, Inches(4.9), "Text-Based Extraction Constraints", accent_color=COLOR_ACCENT_ALT)
    tb = slide7.shapes.add_textbox(Inches(1.0) + col_w + gap, Inches(2.4), col_w - Inches(0.4), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    l2 = [
        "Text-based RAG engines struggle to automatically extract and generate complex diagrams, flowcharts, or formulas inside images.",
        "Complex mathematical notation in unformatted text may require LaTeX formatting.",
        "Non-standard PPTX slide layouts can alter section boundaries."
    ]
    for item in l2:
        p = tf.add_paragraph(); p.text = f"• {item}"; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN; p.space_after = Pt(10)

    create_card(slide7, Inches(0.8) + (col_w + gap) * 2, Inches(1.7), col_w, Inches(4.9), "Ambiguity & API Boundaries", accent_color=COLOR_ACCENT)
    tb = slide7.shapes.add_textbox(Inches(1.0) + (col_w + gap) * 2, Inches(2.4), col_w - Inches(0.4), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    l3 = [
        "Highly ambiguous student queries can retrieve low-relevance vector chunks.",
        "LLM API rate limits (HTTP 429) can occur if queries are submitted in rapid automated succession.",
        "Language limitations exist for highly technical non-English syllabus documents."
    ]
    for item in l3:
        p = tf.add_paragraph(); p.text = f"• {item}"; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN; p.space_after = Pt(10)

    add_footer(slide7, 7)

    # ==========================================
    # SLIDE 8: Future Scope
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8)
    add_header(slide8, "Future Scope & Enhancements")

    create_card(slide8, Inches(0.8), Inches(1.7), card_w2, card_h2, "Multimodal RAG Engine (Vision & Diagrams)", accent_color=COLOR_ACCENT)
    tb = slide8.shapes.add_textbox(Inches(1.0), Inches(2.3), card_w2 - Inches(0.4), card_h2 - Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    p.text = "Integrate Multimodal Gemini Vision models to extract and generate architectural diagrams, flowcharts, circuits, and handwritten lecture notes."
    p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN

    create_card(slide8, Inches(6.933), Inches(1.7), card_w2, card_h2, "Student Progress & Analytics Dashboard", accent_color=COLOR_ACCENT_ALT)
    tb = slide8.shapes.add_textbox(Inches(7.133), Inches(2.3), card_w2 - Inches(0.4), card_h2 - Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    p.text = "Add student performance tracking, quiz score analytics, topic mastery progress bars, and weak-area diagnostic recommendations."
    p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN

    create_card(slide8, Inches(0.8), Inches(4.3), card_w2, card_h2, "Institutional LMS Integration", accent_color=COLOR_ACCENT_ALT)
    tb = slide8.shapes.add_textbox(Inches(1.0), Inches(4.9), card_w2 - Inches(0.4), card_h2 - Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    p.text = "Connect directly with Learning Management Systems (Canvas, Moodle, Blackboard) to auto-sync course syllabi, assignments, and exam dates."
    p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN

    create_card(slide8, Inches(6.933), Inches(4.3), card_w2, card_h2, "Adaptive Learning & Multi-Language Support", accent_color=COLOR_ACCENT)
    tb = slide8.shapes.add_textbox(Inches(7.133), Inches(4.9), card_w2 - Inches(0.4), card_h2 - Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    p.text = "Implement adaptive difficulty quiz generation and real-time multi-lingual translation for non-native English speakers across global curricula."
    p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_MAIN

    add_footer(slide8, 8)

    # ==========================================
    # SLIDE 9: Technology Stack Summary
    # ==========================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide9)
    add_header(slide9, "Technology Stack Summary")

    create_card(slide9, Inches(0.8), Inches(1.7), Inches(11.733), Inches(4.9), "Full Technology Stack Matrix", accent_color=COLOR_ACCENT)

    rows, cols = 5, 4
    table_shape = slide9.shapes.add_table(rows, cols, Inches(1.0), Inches(2.4), Inches(11.333), Inches(3.9))
    table = table_shape.table

    headers = ["Layer", "Technology Used", "Role in Architecture", "Key Highlights"]
    for idx, text in enumerate(headers):
        cell = table.cell(0, idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_ACCENT
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

    data = [
        ("Frontend UI", "Streamlit 1.61+", "User Interface & Tab Navigation", "Executive dark charcoal design system, Google Fonts (Plus Jakarta Sans & Inter)"),
        ("Document Parsers", "pypdf, python-docx, python-pptx", "Multi-Format File Ingestion", "Extracts clean text from PDF, DOCX, PPTX, & TXT documents with overwrite safety"),
        ("Vector Database", "TF-IDF + Cosine Similarity", "Embedding Index & Matcher", "Scikit-learn Vector Space Model with sublinear TF scaling and source-level filtering"),
        ("LLM Backend", "Google Gemini API (gemini-flash-latest)", "RAG Generation & Tutoring", "Strict syllabus-grounded prompts, automatic failover candidates, 4,000+ token limits")
    ]

    for row_idx, row_data in enumerate(data, 1):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(22, 27, 34) if row_idx % 2 == 1 else RGBColor(13, 17, 23)
            p = cell.text_frame.paragraphs[0]
            p.text = cell_value
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT_MAIN

    add_footer(slide9, 9)

    # ==========================================
    # SLIDE 10: Thank You Slide
    # ==========================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide10)

    # Large Center Card
    create_card(slide10, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1), border_accent=True, accent_color=COLOR_ACCENT)

    t_box10 = slide10.shapes.add_textbox(Inches(1.8), Inches(1.6), Inches(9.733), Inches(4.3))
    tf10 = t_box10.text_frame
    tf10.word_wrap = True

    p10_1 = tf10.paragraphs[0]
    p10_1.text = "CAPSTONE PROJECT PRESENTATION"
    p10_1.font.size = Pt(12)
    p10_1.font.bold = True
    p10_1.font.color.rgb = COLOR_EMERALD_LIGHT
    p10_1.space_after = Pt(20)

    p10_2 = tf10.add_paragraph()
    p10_2.text = "THANK YOU!"
    p10_2.font.size = Pt(40)
    p10_2.font.bold = True
    p10_2.font.color.rgb = COLOR_TEXT_MAIN
    p10_2.space_after = Pt(14)

    p10_3 = tf10.add_paragraph()
    p10_3.text = "Questions & Discussion"
    p10_3.font.size = Pt(20)
    p10_3.font.bold = True
    p10_3.font.color.rgb = COLOR_ACCENT_ALT
    p10_3.space_after = Pt(25)

    p10_4 = tf10.add_paragraph()
    p10_4.text = "AI-Based Personalized Tutor Using RAG for Syllabus-Aligned Learning"
    p10_4.font.size = Pt(14)
    p10_4.font.color.rgb = COLOR_TEXT_MUTED

    add_footer(slide10, 10, total_slides=10)

    prs.save(output_path)
    print(f"Successfully generated 10-slide presentation PPTX: {output_path}")

if __name__ == "__main__":
    create_presentation()
